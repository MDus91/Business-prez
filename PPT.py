

from pptx import Presentation
from pptx.util import Inches
import os

# Create a new presentation
prs = Presentation()

# Set paths to your image files (adjust as needed)
logo_path = "logo.png"  # Path to your company logo
placeholder_image = "placeholder.jpg"  # Generic image for each slide

# Define slide content
slides_content = [
    {
        "title": "W-Bank Paint – Protecting Mobile Banking from External Threats",
        "content": "A Modern Approach to App-Level Security\nPresented by: [Your Name], Account Executive, Wultra"
    },
    {
        "title": "The Mobile Banking Challenge",
        "content": "- Mobile banking is growing — and so are threats.\n- Traditional defenses don't protect the app on the device.\n- W-Bank is committed to proactive, real-time protection."
    },
    {
        "title": "The First Layer: Threat Visibility",
        "content": "- Detect compromised devices (rooted/jailbroken)\n- Spot malicious apps and overlay attacks\n- Identify phishing behavior before it strikes"
    },
    {
        "title": "In-App, Real-Time Defense",
        "content": "- Lightweight SDK runs silently inside the app\n- Responds to threats instantly:\n  - Disable functions\n  - Notify users\n  - Block access if needed"
    },
    {
        "title": "Continuous Threat Intelligence",
        "content": "- Updated malware signatures\n- Cloud-supported detection patterns\n- Compliant with PSD2, GDPR"
    },
    {
        "title": "Customer Confidence as a Security Outcome",
        "content": "- Transparent security builds user trust\n- Visible protection reassures app users\n- Reduces fraud and reputational risk"
    },
    {
        "title": "Why W-Bank Paint Works",
        "content": "- Embedded, real-time threat defense\n- Smart, responsive in-app controls\n- Adaptive protection against new threats\n- Trusted by financial institutions"
    },
    {
        "title": "Ready to Secure Your App",
        "content": "- Pilot available for testing in your environment\n- Fast, guided integration\n- Full support from Wultra’s mobile security team"
    }
]

# Add slides and content
for slide_info in slides_content:
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = slide_info["title"]
    content.text = slide_info["content"]

    # Add logo (bottom right)
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(9), Inches(6.7), height=Inches(0.5))

    # Add image (left side)
    if os.path.exists(placeholder_image):
        slide.shapes.add_picture(placeholder_image, Inches(0.5), Inches(3.5), width=Inches(3.5))

# Save the PowerPoint
prs.save("W-Bank_Paint_Presentation_With_Images.pptx")
